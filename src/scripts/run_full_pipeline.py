import io
import sys
from pathlib import Path
from datetime import datetime
import joblib
import pandas as pd
import streamlit as st
import plotly.express as px
from sklearn.feature_extraction.text import TfidfVectorizer, CountVectorizer
from sklearn.feature_extraction import text as sk_text
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score
from concurrent.futures import ThreadPoolExecutor, as_completed, TimeoutError
import json
import time
import ollama
import ast
import re

# PDF export imports
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    HRFlowable, PageBreak, Image as RLImage,
)

PROJECT_ROOT    = Path(__file__).resolve().parents[2]
DATA_DIR        = PROJECT_ROOT / "src" / "sample_data"
OUTPUT_DIR      = PROJECT_ROOT / "output"
MODEL_FILE      = OUTPUT_DIR / "sentiment_model.pkl"
VECTORIZER_FILE = OUTPUT_DIR / "tfidf_vectorizer.pkl"

THEMES = [
    "Product Quality",
    "Product Availability",
    "Customer Service",
    "Speed of Service",
    "Store Environment",
    "Price & Value",
    "Digital & Rewards",
    "Policies & Safety",
]

st.set_page_config(page_title="Feedback Intelligence Platform", layout="wide")


# ── Shared helpers ─────────────────────────────────────────────────────────────

COLOUR_MAP = {
    "positive":      "#16a34a",
    "negative":      "#dc2626",
    "neutral":       "#6b7280",
    "neutral/mixed": "#9ca3af",
}

def _text_col(df):
    for c in ["text", "raw_text", "clean_text"]:
        if c in df.columns:
            return c
    return "clean_text"


# ── LLM ───────────────────────────────────────────────────────────────────────

def call_llm(prompt, model="gemma2:9b"):
    try:
        response = ollama.chat(model=model, messages=[{"role": "user", "content": prompt}])
        return response["message"]["content"]
    except Exception as e:
        print(f"Error calling LLM: {e}")
        raise e


# ── Model training / loading ──────────────────────────────────────────────────

def ensure_output_dir():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def sentiments_from_stars(stars, classification_type="three_class"):
    try:
        stars = float(stars)
    except (TypeError, ValueError):
        return None
    if classification_type == "binary":
        return "positive" if stars >= 4 else ("negative" if stars <= 2 else None)
    return "positive" if stars >= 4 else ("neutral/mixed" if stars == 3 else "negative")


def prepare_training_data():
    train_file = DATA_DIR / "training_testing_data.csv"
    if not train_file.exists():
        candidates = sorted(DATA_DIR.glob("*.csv"))
        if candidates:
            train_file = candidates[0]
        else:
            raise FileNotFoundError("No training CSV found in src/sample_data")
    df = pd.read_csv(train_file)
    if "sentiment" not in df.columns:
        df["sentiment"] = df["stars"].apply(lambda x: sentiments_from_stars(x, "three_class"))
    if "clean_text" not in df.columns:
        if "text" in df.columns:
            df["clean_text"] = df["text"].fillna("").astype(str).str.lower()
        elif "raw_text" in df.columns:
            df["clean_text"] = df["raw_text"].fillna("").astype(str).str.lower()
        else:
            raise ValueError("Training data must contain 'clean_text', 'text', or 'raw_text'")
    return df[df["sentiment"].notna()].copy()


def train_model():
    df = prepare_training_data()
    content, sent = df["clean_text"], df["sentiment"]
    c_train, c_test, s_train, s_test = train_test_split(
        content, sent, test_size=0.2, random_state=2016, stratify=sent)
    extra_stop = {"review", "user", "star", "stars", "https", "http", "amp"}
    stop_words  = set(sk_text.ENGLISH_STOP_WORDS) | extra_stop
    vectorizer  = TfidfVectorizer(max_features=5000, ngram_range=(1, 2),
                                   min_df=2, max_df=0.8, stop_words=list(stop_words))
    X_train = vectorizer.fit_transform(c_train)
    X_test  = vectorizer.transform(c_test)
    model   = LogisticRegression(max_iter=1000, random_state=2016, C=0.8)
    model.fit(X_train, s_train)
    accuracy = accuracy_score(s_test, model.predict(X_test))
    joblib.dump(model, MODEL_FILE)
    joblib.dump(vectorizer, VECTORIZER_FILE)
    return model, vectorizer, accuracy


def load_or_train_model():
    ensure_output_dir()
    if MODEL_FILE.exists() and VECTORIZER_FILE.exists():
        return joblib.load(MODEL_FILE), joblib.load(VECTORIZER_FILE), None
    with st.spinner("Training model from sample data..."):
        return train_model()


# ── Preprocessing & prediction ────────────────────────────────────────────────

def preprocess_reviews(df):
    if "clean_text" in df.columns:
        df["clean_text"] = df["clean_text"].fillna("").astype(str).str.lower()
    elif "text" in df.columns:
        df["clean_text"] = df["text"].fillna("").astype(str).str.lower()
    elif "raw_text" in df.columns:
        df["clean_text"] = df["raw_text"].fillna("").astype(str).str.lower()
    else:
        raise ValueError("CSV must have a 'text', 'raw_text', or 'clean_text' column")
    return df


CONTRAST_WORDS = {"but", "however", "though", "although", "yet", "except", "overall", "while"}
POS_CUES = {"good", "great", "nice", "friendly", "fast", "clean", "love", "excellent", "amazing", "enjoy"}
NEG_CUES = {"bad", "slow", "rude", "wrong", "dirty", "hate", "awful", "terrible", "issue", "problem"}


def has_contrast(text):
    t = f" {text.lower()} "
    return any(f" {w} " in t for w in CONTRAST_WORDS)


def has_dual_polarity_words(text):
    tokens = set(re.findall(r"[a-z']+", text.lower()))
    return bool(tokens & POS_CUES) and bool(tokens & NEG_CUES)


def mixed_rule(row):
    text  = str(row.get("clean_text", ""))
    p_pos = float(row.get("prob_positive", 0.0))
    p_neg = float(row.get("prob_negative", 0.0))
    prob_cond     = (p_pos >= 0.30) and (p_neg >= 0.30) and (abs(p_pos - p_neg) <= 0.25)
    contrast_cond = has_contrast(text)
    lex_cond      = has_dual_polarity_words(text)
    return (prob_cond and contrast_cond) or (contrast_cond and lex_cond)


def predict_reviews(df, model, vectorizer):
    tfidf = vectorizer.transform(df["clean_text"])
    preds = model.predict(tfidf)
    probs = model.predict_proba(tfidf)
    df["predicted_sentiment"] = preds
    df["confidence"] = probs.max(axis=1)
    for idx, cls in enumerate(model.classes_):
        df[f"prob_{cls}"] = probs[:, idx]
    df["is_mixed"] = False
    mask = df["predicted_sentiment"].isin(["neutral", "neutral/mixed"])
    if mask.any():
        df.loc[mask, "is_mixed"] = df[mask].apply(mixed_rule, axis=1)
    return df


# ── Theme extraction ──────────────────────────────────────────────────────────

def build_prompt(batch, themes):
    numbered = "\n".join([f"Review {i+1}:\n{str(r)[:250]}" for i, r in enumerate(batch)])
    return f"""You are a professional theme classifier for customer reviews.
You work for a Feedback Intelligence Platform that analyzes reviews.

Available themes: {themes}

RULES:
- Only assign themes from the available list above. NEVER invent or create your own themes.
- Every review MUST have at least one theme assigned.
- Return ONLY a valid JSON dictionary where keys are Review Numbers ("1", "2", etc.) and values are arrays of themes.
- You MUST generate exactly {len(batch)} keys.
- No extra explanation outside the JSON block.

Example for 3 reviews:
{{
  "1": ["Customer Service", "Product Quality"],
  "2": ["Speed of Service"],
  "3": ["Price & Value"]
}}

Reviews to classify:
{numbered}

Return ONLY valid JSON.
"""


def extract_themes_with_retry(batch_info, themes_list, max_retries=5):
    batch_idx, batch = batch_info
    prompt = build_prompt(batch, themes_list)
    for attempt in range(1, max_retries + 1):
        try:
            raw         = call_llm(prompt)
            first_brace = raw.find("{")
            first_brack = raw.find("[")
            start = min(i for i in [first_brace, first_brack] if i != -1) \
                    if any(i != -1 for i in [first_brace, first_brack]) else -1
            end   = max(raw.rfind("}"), raw.rfind("]")) + 1
            if start == -1 or end <= start:
                raise ValueError("LLM output contained no valid JSON")
            clean = raw[start:end]
            try:
                parsed_data = json.loads(clean)
            except json.JSONDecodeError:
                try:
                    parsed_data = ast.literal_eval(clean)
                except Exception:
                    if clean.strip().startswith("{") and clean.strip().endswith("}"):
                        parsed_data = json.loads(f"[{clean}]")
                    else:
                        raise ValueError(f"Could not parse: {clean}")
            parsed_dict = {}
            if isinstance(parsed_data, list):
                for item in parsed_data:
                    if isinstance(item, dict):
                        parsed_dict.update(item)
            elif isinstance(parsed_data, dict):
                parsed_dict = parsed_data
            else:
                raise ValueError("Response is not valid JSON")
            themes = []
            for idx in range(1, len(batch) + 1):
                t = parsed_dict.get(str(idx), [])
                if not isinstance(t, list):
                    t = [t]
                if not t:
                    t = ["Customer Service"]
                themes.append(t)
            if len(themes) != len(batch):
                raise ValueError(f"Count mismatch: got {len(themes)}, expected {len(batch)}")
            validated = []
            for tlist in themes:
                safe  = [str(list(t.values())[0]) if isinstance(t, dict) and t else str(t) for t in tlist]
                valid = []
                for st_t in safe:
                    for real in themes_list:
                        if st_t.strip().lower() == real.lower():
                            if real not in valid:
                                valid.append(real)
                            break
                if not valid:
                    raise ValueError(f"Hallucination: '{tlist}' not in allowed list")
                validated.append(valid)
            return batch_idx, batch, validated, "success"
        except Exception as e:
            print(f"Batch {batch_idx} attempt {attempt} failed: {e}")
            if attempt < max_retries:
                time.sleep(2)
    return batch_idx, batch, None, "failed"


def extract_themes(df, themes_list, batch_size=30, max_workers=2):
    reviews = df["clean_text"].fillna("").tolist()
    batches = [(i, reviews[i:i+batch_size]) for i in range(0, len(reviews), batch_size)]
    success, failed = [], []
    pbar   = st.progress(0)
    status = st.empty()
    total  = len(batches)
    done   = 0
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {ex.submit(extract_themes_with_retry, b, themes_list): b for b in batches}
        for future in as_completed(futures):
            try:
                bidx, batch, bthemes, bstatus = future.result(timeout=180)
            except TimeoutError:
                bidx, batch = futures[future]
                bstatus, bthemes = "failed", [["FAILED"]] * len(batch)
            except Exception:
                bidx, batch = futures[future]
                bstatus, bthemes = "failed", [["FAILED"]] * len(batch)
            if bstatus == "success":
                for i, (_, vt) in enumerate(zip(batch, bthemes)):
                    success.append({"original_idx": bidx + i, "themes": ", ".join(vt)})
            else:
                for i in range(len(batch)):
                    failed.append({"original_idx": bidx + i, "themes": "FAILED"})
            done += 1
            if done % max(1, total // 100) == 0 or done == total:
                pbar.progress(done / total)
                status.text(f"Processed batch {done}/{total} — please wait...")
    lookup = {r["original_idx"]: r["themes"] for r in success + failed}
    df["themes"] = [lookup.get(i, "FAILED") for i in range(len(df))]
    before = len(df)
    df = df[~df["themes"].str.contains("FAILED", na=False)]
    if len(df) < before:
        print(f"Dropped {before - len(df)} reviews due to extraction failure.")
    pbar.progress(1.0)
    status.text("Theme extraction complete!")
    return df


# ── PDF export ────────────────────────────────────────────────────────────────

_CLR_POS    = colors.HexColor("#16a34a")
_CLR_NEG    = colors.HexColor("#dc2626")
_CLR_HEADER = colors.HexColor("#1e3a5f")
_CLR_RULE   = colors.HexColor("#e5e7eb")
_CLR_BG     = colors.HexColor("#f0f4ff")


def _build_pdf_styles():
    base = getSampleStyleSheet()
    return (
        ParagraphStyle("T",  parent=base["Title"],    fontSize=22, textColor=_CLR_HEADER, spaceAfter=4,   fontName="Helvetica-Bold"),
        ParagraphStyle("ST", parent=base["Normal"],   fontSize=10, textColor=colors.HexColor("#6b7280"),  spaceAfter=16),
        ParagraphStyle("H2", parent=base["Heading2"], fontSize=13, textColor=_CLR_HEADER, spaceBefore=18, spaceAfter=6, fontName="Helvetica-Bold"),
        ParagraphStyle("B",  parent=base["Normal"],   fontSize=9,  leading=14, textColor=colors.HexColor("#374151")),
        ParagraphStyle("S",  parent=base["Normal"],   fontSize=8,  leading=12, textColor=colors.HexColor("#6b7280")),
    )


def _fig_to_img(fig, w=6.5, h=3.2):
    buf = io.BytesIO(fig.to_image(format="png", width=int(w*100), height=int(h*100), scale=2))
    return RLImage(buf, width=w*inch, height=h*inch)


def _pdf_pie(df):
    counts = df["predicted_sentiment"].value_counts().reset_index()
    counts.columns = ["sentiment", "count"]
    fig = px.pie(counts, names="sentiment", values="count",
                 color="sentiment", color_discrete_map=COLOUR_MAP, title="Sentiment Distribution")
    fig.update_layout(margin=dict(l=20, r=20, t=40, b=20), paper_bgcolor="white", font_size=11)
    return _fig_to_img(fig, 3.0, 2.8)


def _pdf_bar(df):
    exp = df["themes"].str.split(r",\s*").explode().str.strip()
    exp = exp[~exp.isin(["FAILED", "", "NOT PROCESSED"])].value_counts().head(8).reset_index()
    exp.columns = ["Theme", "Count"]
    fig = px.bar(exp, x="Count", y="Theme", orientation="h", title="Top Themes",
                 color="Count", color_continuous_scale="Blues")
    fig.update_layout(margin=dict(l=20, r=20, t=40, b=20), paper_bgcolor="white",
                      yaxis=dict(categoryorder="total ascending"), coloraxis_showscale=False)
    return _fig_to_img(fig, 3.8, 2.8)


def _pdf_trend(df):
    if "date" not in df.columns:
        return None
    try:
        d  = df.copy()
        d["date"] = pd.to_datetime(d["date"])
        td = d.groupby([pd.Grouper(key="date", freq="ME"), "predicted_sentiment"])\
               .size().reset_index(name="count")
        if len(td) <= 1:
            return None
        fig = px.line(td, x="date", y="count", color="predicted_sentiment",
                      title="Monthly Trend", color_discrete_map=COLOUR_MAP)
        fig.update_layout(margin=dict(l=20, r=20, t=40, b=20), paper_bgcolor="white")
        return _fig_to_img(fig, 6.5, 2.8)
    except Exception:
        return None


def _pdf_kpi(df):
    total   = len(df)
    pos     = int((df["predicted_sentiment"] == "positive").sum())
    neg     = int((df["predicted_sentiment"] == "negative").sum())
    mixed   = int(df["predicted_sentiment"].isin(["neutral", "neutral/mixed"]).sum())
    conf    = df["confidence"].mean() if "confidence" in df.columns else None
    flagged = int((df["is_mixed"] == True).sum()) if "is_mixed" in df.columns else 0
    rows = [
        ["Metric", "Value", "Share"],
        ["Total reviews analysed",  f"{total:,}",    "100%"],
        ["Positive",                f"{pos:,}",       f"{pos/total*100:.1f}%"],
        ["Negative",                f"{neg:,}",       f"{neg/total*100:.1f}%"],
        ["Neutral / Mixed",         f"{mixed:,}",     f"{mixed/total*100:.1f}%"],
        ["Mixed-signal flagged",    f"{flagged:,}",   f"{flagged/total*100:.1f}%"],
    ]
    if conf:
        rows.append(["Avg model confidence", f"{conf*100:.1f}%", "—"])
    tbl = Table(rows, colWidths=[3.2*inch, 1.5*inch, 1.5*inch])
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0), _CLR_BG),
        ("TEXTCOLOR",     (0, 0), (-1, 0), _CLR_HEADER),
        ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, 0), 9),
        ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 1), (-1, -1), 9),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
        ("TEXTCOLOR",     (0, 1), (-1, -1), colors.HexColor("#374151")),
        ("TEXTCOLOR",     (2, 2), (2, 2),  _CLR_POS),
        ("TEXTCOLOR",     (2, 3), (2, 3),  _CLR_NEG),
        ("GRID",          (0, 0), (-1, -1), 0.4, _CLR_RULE),
        ("TOPPADDING",    (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING",   (0, 0), (-1, -1), 8),
    ]))
    return tbl


def _pdf_theme_table(df):
    if "themes" not in df.columns:
        return None
    exp = df.assign(Theme=df["themes"].str.split(r",\s*")).explode("Theme")
    exp["Theme"] = exp["Theme"].str.strip()
    exp = exp[~exp["Theme"].isin(["FAILED", "", "NOT PROCESSED"])]
    if exp.empty:
        return None
    pivot = pd.crosstab(exp["Theme"], exp["predicted_sentiment"])
    pivot["Total"] = pivot.sum(axis=1)
    pivot = pivot.sort_values("Total", ascending=False).head(8)
    scols = [c for c in ["positive", "negative", "neutral", "neutral/mixed"] if c in pivot.columns]
    rows  = [["Theme"] + [c.capitalize() for c in scols] + ["Total"]]
    for theme, row in pivot.iterrows():
        dr = [str(theme)]
        for sc in scols:
            v = int(row.get(sc, 0))
            p = v / row["Total"] * 100 if row["Total"] else 0
            dr.append(f"{v} ({p:.0f}%)")
        dr.append(str(int(row["Total"])))
        rows.append(dr)
    tbl = Table(rows, colWidths=[2.5*inch] + [1.0*inch]*len(scols) + [0.8*inch])
    style = [
        ("BACKGROUND",    (0, 0), (-1, 0), _CLR_BG),
        ("TEXTCOLOR",     (0, 0), (-1, 0), _CLR_HEADER),
        ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, 0), 8),
        ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 1), (-1, -1), 8),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
        ("GRID",          (0, 0), (-1, -1), 0.4, _CLR_RULE),
        ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
        ("ALIGN",         (1, 0), (-1, -1), "CENTER"),
    ]
    for i, sc in enumerate(scols):
        if sc == "positive": style.append(("TEXTCOLOR", (i+1, 0), (i+1, 0), _CLR_POS))
        if sc == "negative": style.append(("TEXTCOLOR", (i+1, 0), (i+1, 0), _CLR_NEG))
    tbl.setStyle(TableStyle(style))
    return tbl


def _pdf_samples(df):
    tc = _text_col(df)
    get = lambda s: (df[df["predicted_sentiment"] == s].nlargest(3, "confidence")
                     if "confidence" in df.columns
                     else df[df["predicted_sentiment"] == s].head(3))
    sample = pd.concat([get("positive"), get("negative")])
    rows = [["Sentiment", "Review excerpt", "Confidence"]]
    for _, r in sample.iterrows():
        ex   = str(r[tc])[:160].strip() + ("…" if len(str(r[tc])) > 160 else "")
        conf = f"{r['confidence']*100:.0f}%" if "confidence" in r else "—"
        rows.append([r["predicted_sentiment"].capitalize(), ex, conf])
    tbl = Table(rows, colWidths=[1.0*inch, 4.8*inch, 0.9*inch])
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, 0), _CLR_BG),
        ("TEXTCOLOR",     (0, 0), (-1, 0), _CLR_HEADER),
        ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, 0), 8),
        ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 1), (-1, -1), 8),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
        ("GRID",          (0, 0), (-1, -1), 0.4, _CLR_RULE),
        ("TOPPADDING",    (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
    ]))
    return tbl


def _pdf_footer(canvas, doc):
    canvas.saveState()
    canvas.setFont("Helvetica", 7)
    canvas.setFillColor(colors.HexColor("#9ca3af"))
    canvas.drawString(0.75*inch, 0.5*inch,
                      "Customer Feedback Intelligence Platform  |  Confidential")
    canvas.drawRightString(letter[0]-0.75*inch, 0.5*inch, f"Page {doc.page}")
    canvas.restoreState()


def generate_pdf_report(df, report_title="Feedback Intelligence Report"):
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter,
                            leftMargin=0.75*inch, rightMargin=0.75*inch,
                            topMargin=0.85*inch, bottomMargin=0.75*inch)
    ts, sts, h2s, bs, ss = _build_pdf_styles()
    story = []
    story.append(Paragraph(report_title, ts))
    story.append(Paragraph(
        f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}  |  "
        f"{len(df):,} reviews analysed", sts))
    story.append(HRFlowable(width="100%", thickness=1, color=_CLR_RULE, spaceAfter=12))
    story.append(Paragraph("1. Summary Metrics", h2s))
    story.append(_pdf_kpi(df))
    story.append(Spacer(1, 14))
    story.append(Paragraph("2. Sentiment & Theme Overview", h2s))
    try:
        pie = _pdf_pie(df)
        bar = _pdf_bar(df) if "themes" in df.columns else None
        if bar:
            cr = Table([[pie, bar]], colWidths=[3.2*inch, 4.0*inch])
            cr.setStyle(TableStyle([("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                                    ("LEFTPADDING", (0,0), (-1,-1), 0),
                                    ("RIGHTPADDING", (0,0), (-1,-1), 8)]))
            story.append(cr)
        else:
            story.append(pie)
    except Exception as e:
        story.append(Paragraph(f"Charts unavailable: {e} — ensure kaleido is installed.", ss))
    story.append(Spacer(1, 10))
    trend = _pdf_trend(df)
    off   = 0
    if trend:
        story.append(Paragraph("3. Monthly Sentiment Trend", h2s))
        story.append(trend)
        story.append(Spacer(1, 10))
        off = 1
    tt = _pdf_theme_table(df)
    if tt:
        story.append(Paragraph(f"{3+off}. Theme Sentiment Breakdown", h2s))
        story.append(Paragraph("Count (%) of reviews per sentiment for the top 8 themes.", bs))
        story.append(Spacer(1, 6))
        story.append(tt)
        story.append(Spacer(1, 10))
    st_ = _pdf_samples(df)
    if st_:
        story.append(PageBreak())
        story.append(Paragraph(f"{4+off}. Sample Review Verbatims", h2s))
        story.append(Paragraph("Highest-confidence positive and negative reviews.", bs))
        story.append(Spacer(1, 6))
        story.append(st_)
    story.append(Spacer(1, 20))
    story.append(HRFlowable(width="100%", thickness=0.5, color=_CLR_RULE, spaceAfter=6))
    story.append(Paragraph(
        "Auto-generated by the Customer Feedback Intelligence Platform. "
        "Sentiment predictions use Logistic Regression. Themes assigned by local LLM (Ollama/Gemma). "
        "Review alongside raw data before making decisions.", ss))
    doc.build(story, onFirstPage=_pdf_footer, onLaterPages=_pdf_footer)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def _fig_to_png_bytes(fig, width=900, height=500):
    """Convert a Plotly figure to PNG bytes for embedding in PowerPoint."""
    return fig.to_image(format="png", width=width, height=height, scale=2)


def generate_pptx(df: pd.DataFrame, report_title: str = "Feedback Intelligence Report") -> bytes:
    """
    Build a PowerPoint deck from the analysed dataframe.
    Slides:
        1. Title slide
        2. Key metrics
        3. Sentiment distribution pie chart
        4. Star rating vs sentiment (if stars column exists)
        5. Top themes bar chart
        6. Theme sentiment breakdown table
        7. Positive reviews — top themes
        8. Negative reviews — top themes
        9. Key findings & recommendations
    Returns raw .pptx bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Colour palette ─────────────────────────────────────────────────────────
    NAVY    = RGBColor(0x1e, 0x3a, 0x5f)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN   = RGBColor(0x16, 0xa3, 0x4a)
    RED     = RGBColor(0xdc, 0x26, 0x26)
    GRAY    = RGBColor(0x6b, 0x72, 0x80)
    LIGHT   = RGBColor(0xf8, 0xf9, 0xfa)
    ACCENT  = RGBColor(0x06, 0x5a, 0x82)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLANK = prs.slide_layouts[6]   # completely blank layout

    # ── Helper functions ───────────────────────────────────────────────────────

    def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.line.fill.background()
        if alpha is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            shape.fill.fore_color.transparency = alpha
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        return shape

    def add_text(slide, text, x, y, w, h, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        txBox.word_wrap = wrap
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color else NAVY
        return txBox

    def add_chart_image(slide, fig, x, y, w, h, width_px=900, height_px=500):
        try:
            png = _fig_to_png_bytes(fig, width=width_px, height=height_px)
            img_stream = io.BytesIO(png)
            slide.shapes.add_picture(img_stream, Inches(x), Inches(y),
                                     Inches(w), Inches(h))
        except Exception:
            add_text(slide, "Chart unavailable — install kaleido",
                     x, y, w, h, size=12, color=GRAY)

    def slide_header(slide, title, subtitle=""):
        """Navy top bar with white title text."""
        add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
        add_text(slide, title, 0.4, 0.1, 11, 0.75,
                 size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        if subtitle:
            add_text(slide, subtitle, 0.4, 0.78, 11, 0.35,
                     size=13, bold=False, color=RGBColor(0xCA, 0xDC, 0xFC),
                     align=PP_ALIGN.LEFT)

    # ── Pre-compute stats ──────────────────────────────────────────────────────
    total    = len(df)
    pos      = int((df["predicted_sentiment"] == "positive").sum())
    neg      = int((df["predicted_sentiment"] == "negative").sum())
    mixed    = int(df["predicted_sentiment"].isin(["neutral", "neutral/mixed"]).sum())
    avg_conf = df["confidence"].mean() * 100 if "confidence" in df.columns else None
    has_themes = "themes" in df.columns
    has_stars  = "stars" in df.columns

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(BLANK)
    add_rect(s1, 0, 0, 13.33, 7.5, NAVY)
    add_rect(s1, 0, 4.8, 13.33, 2.7, ACCENT)
    add_text(s1, report_title,
             0.6, 1.6, 12, 1.4, size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s1, "Customer Feedback Intelligence Platform",
             0.6, 3.1, 12, 0.6, size=18, bold=False,
             color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.LEFT)
    add_text(s1, f"Generated {datetime.now().strftime('%B %d, %Y')}  |  {total:,} reviews analysed",
             0.6, 3.7, 12, 0.5, size=14, bold=False,
             color=RGBColor(0x9F, 0xB3, 0xD8), align=PP_ALIGN.LEFT)
    add_text(s1, "Sentiment · Themes · Insights",
             0.6, 5.2, 12, 0.6, size=16, bold=False,
             color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Key Metrics
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(BLANK)
    slide_header(s2, "Key Metrics", f"Based on {total:,} customer reviews")
    add_rect(s2, 0, 1.1, 13.33, 6.4, LIGHT)

    metrics = [
        ("Total Reviews",    f"{total:,}",              NAVY,  "All reviews analysed"),
        ("Positive",         f"{pos:,}",                GREEN, f"{pos/total*100:.1f}% of total"),
        ("Negative",         f"{neg:,}",                RED,   f"{neg/total*100:.1f}% of total"),
        ("Neutral / Mixed",  f"{mixed:,}",              GRAY,  f"{mixed/total*100:.1f}% of total"),
        ("Avg Confidence",   f"{avg_conf:.1f}%" if avg_conf else "N/A",
                                                        ACCENT,"Model certainty"),
    ]

    card_w = 2.3
    card_gap = 0.26
    start_x = 0.35

    for i, (label, value, clr, sub) in enumerate(metrics):
        cx = start_x + i * (card_w + card_gap)
        # card background
        card = s2.shapes.add_shape(1, Inches(cx), Inches(1.5),
                                    Inches(card_w), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
        card.line.width = Emu(9525)   # 0.75pt
        # coloured top accent bar
        top = s2.shapes.add_shape(1, Inches(cx), Inches(1.5),
                                   Inches(card_w), Inches(0.12))
        top.fill.solid()
        top.fill.fore_color.rgb = clr
        top.line.fill.background()
        # label
        add_text(s2, label, cx + 0.15, 1.75, card_w - 0.3, 0.45,
                 size=12, bold=False, color=GRAY, align=PP_ALIGN.LEFT)
        # value
        add_text(s2, value, cx + 0.15, 2.25, card_w - 0.3, 1.1,
                 size=34, bold=True, color=clr, align=PP_ALIGN.LEFT)
        # sub-label
        add_text(s2, sub, cx + 0.15, 3.45, card_w - 0.3, 0.6,
                 size=11, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Sentiment Distribution Pie
    # ══════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(BLANK)
    slide_header(s3, "Sentiment Distribution", "Overall breakdown of predicted sentiment")
    add_rect(s3, 0, 1.1, 13.33, 6.4, LIGHT)

    counts = df["predicted_sentiment"].value_counts().reset_index()
    counts.columns = ["sentiment", "count"]
    fig_pie = px.pie(counts, names="sentiment", values="count",
                     color="sentiment", color_discrete_map=COLOUR_MAP,
                     title="")
    fig_pie.update_layout(margin=dict(l=10, r=10, t=10, b=10),
                          paper_bgcolor="white", showlegend=True,
                          legend=dict(font=dict(size=14)))
    add_chart_image(s3, fig_pie, 1.5, 1.4, 6.5, 5.5, width_px=700, height_px=560)

    # Stat callouts on the right
    callouts = [
        (f"{pos:,}", f"Positive  ({pos/total*100:.0f}%)", GREEN),
        (f"{neg:,}", f"Negative  ({neg/total*100:.0f}%)", RED),
        (f"{mixed:,}", f"Neutral/Mixed  ({mixed/total*100:.0f}%)", GRAY),
    ]
    for i, (val, lbl, clr) in enumerate(callouts):
        cy = 2.0 + i * 1.7
        add_rect(s3, 8.5, cy, 4.2, 1.35, WHITE)
        add_text(s3, val, 8.7, cy + 0.05, 3.8, 0.7,
                 size=30, bold=True, color=clr, align=PP_ALIGN.LEFT)
        add_text(s3, lbl, 8.7, cy + 0.75, 3.8, 0.5,
                 size=12, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Star Rating vs Sentiment (conditional)
    # ══════════════════════════════════════════════════════════════════════════
    if has_stars:
        s4 = prs.slides.add_slide(BLANK)
        slide_header(s4, "Star Rating vs Model Prediction",
                     "How well does model sentiment agree with reviewer star ratings?")
        add_rect(s4, 0, 1.1, 13.33, 6.4, LIGHT)
        try:
            star_df = df.copy()
            star_df["stars"] = pd.to_numeric(star_df["stars"], errors="coerce")
            star_df = star_df.dropna(subset=["stars"])
            star_df["stars"] = star_df["stars"].astype(int)
            star_df = star_df[star_df["stars"].between(1, 5)]
            star_counts = (star_df.groupby(["stars", "predicted_sentiment"])
                           .size().reset_index(name="count"))
            fig_stars = px.bar(star_counts, x="stars", y="count",
                               color="predicted_sentiment",
                               color_discrete_map=COLOUR_MAP, barmode="stack",
                               labels={"stars": "Star Rating", "count": "Reviews"})
            fig_stars.update_traces(marker_line_color="white", marker_line_width=1)
            fig_stars.update_layout(
                bargap=0.2, paper_bgcolor="white",
                xaxis=dict(tickmode="array", tickvals=[1,2,3,4,5],
                           ticktext=["1★","2★","3★","4★","5★"]),
                yaxis=dict(rangemode="tozero", dtick=1, tickformat="d"),
                legend_title_text="Model Prediction",
                margin=dict(l=20, r=20, t=20, b=20),
            )
            add_chart_image(s4, fig_stars, 0.5, 1.3, 8.5, 5.8, width_px=860, height_px=580)

            # Agreement stats
            def exp_sent(stars):
                return "positive" if stars >= 4 else ("negative" if stars <= 2 else "neutral/mixed")
            star_df["expected"] = star_df["stars"].apply(exp_sent)
            star_df["agrees"]   = (star_df["predicted_sentiment"] == star_df["expected"]) | \
                                  ((star_df["predicted_sentiment"] == "neutral") &
                                   (star_df["expected"] == "neutral/mixed"))
            overall = star_df["agrees"].mean() * 100

            add_rect(s4, 9.3, 1.5, 3.6, 1.4, WHITE)
            add_text(s4, f"{overall:.1f}%", 9.5, 1.55, 3.2, 0.75,
                     size=38, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
            add_text(s4, "Overall model agreement\nwith star ratings",
                     9.5, 2.35, 3.2, 0.55, size=12, color=GRAY)

            add_rect(s4, 9.3, 3.1, 3.6, 3.6, WHITE)
            add_text(s4, "Per-star agreement", 9.5, 3.2, 3.2, 0.4,
                     size=13, bold=True, color=NAVY)
            by_star = star_df.groupby("stars")["agrees"].agg(["sum","count"]).reset_index()
            for idx, row in by_star.iterrows():
                pct = row["sum"] / row["count"] * 100
                ry  = 3.65 + idx * 0.55
                add_text(s4, f"{int(row['stars'])} star", 9.5, ry, 1.5, 0.45, size=11, color=GRAY)
                add_text(s4, f"{pct:.0f}%", 11.0, ry, 1.7, 0.45,
                         size=11, bold=True,
                         color=GREEN if pct >= 70 else (RED if pct < 50 else GRAY))
        except Exception:
            add_text(s4, "Could not generate star chart.", 0.5, 2, 12, 1, size=14, color=GRAY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Top Themes
    # ══════════════════════════════════════════════════════════════════════════
    if has_themes:
        s5 = prs.slides.add_slide(BLANK)
        slide_header(s5, "Top Themes", "Most frequently mentioned topics across all reviews")
        add_rect(s5, 0, 1.1, 13.33, 6.4, LIGHT)

        exp_all = df["themes"].str.split(r",\s*").explode().str.strip()
        exp_all = exp_all[~exp_all.isin(["FAILED", "", "NOT PROCESSED"])]
        tc = exp_all.value_counts().head(8).reset_index()
        tc.columns = ["Theme", "Count"]

        fig_bar = px.bar(tc, x="Count", y="Theme", orientation="h",
                         color="Count", color_continuous_scale="Blues",
                         labels={"Count": "Mentions"})
        fig_bar.update_traces(marker_line_color="white", marker_line_width=1)
        fig_bar.update_layout(
            bargap=0.2, paper_bgcolor="white", coloraxis_showscale=False,
            yaxis=dict(categoryorder="total ascending"),
            margin=dict(l=20, r=20, t=20, b=20),
        )
        add_chart_image(s5, fig_bar, 0.4, 1.3, 8.5, 5.8, width_px=860, height_px=580)

        # Top 3 callouts
        top3 = tc.head(3)
        for i, (_, row) in enumerate(top3.iterrows()):
            cy = 1.5 + i * 1.8
            add_rect(s5, 9.3, cy, 3.6, 1.55, WHITE)
            add_text(s5, f"#{i+1}", 9.5, cy + 0.05, 0.7, 0.5,
                     size=18, bold=True, color=ACCENT)
            add_text(s5, row["Theme"], 10.2, cy + 0.05, 2.6, 0.55,
                     size=13, bold=True, color=NAVY)
            add_text(s5, f"{int(row['Count'])} mentions", 9.5, cy + 0.65, 3.2, 0.4,
                     size=12, color=GRAY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Theme Sentiment Breakdown
    # ══════════════════════════════════════════════════════════════════════════
    if has_themes:
        s6 = prs.slides.add_slide(BLANK)
        slide_header(s6, "Theme Sentiment Breakdown",
                     "Positive / negative / neutral split per theme")
        add_rect(s6, 0, 1.1, 13.33, 6.4, LIGHT)

        df_exp = df.assign(Theme=df["themes"].str.split(r",\s*")).explode("Theme")
        df_exp["Theme"] = df_exp["Theme"].str.strip()
        df_exp = df_exp[~df_exp["Theme"].isin(["FAILED", "", "NOT PROCESSED"])]

        pivot = pd.crosstab(df_exp["Theme"], df_exp["predicted_sentiment"])
        pivot["Total"] = pivot.sum(axis=1)
        pivot = pivot.sort_values("Total", ascending=False).head(6)
        scols = [c for c in ["positive", "negative", "neutral", "neutral/mixed"]
                 if c in pivot.columns]

        # Table headers
        col_headers = ["Theme"] + [c.capitalize() for c in scols] + ["Total"]
        n_cols = len(col_headers)
        col_widths = [3.5] + [1.6] * (n_cols - 2) + [1.2]
        start_x = 0.4
        header_y = 1.3

        # Header row
        cx = start_x
        for j, hdr in enumerate(col_headers):
            cw = col_widths[j]
            rect = s6.shapes.add_shape(1, Inches(cx), Inches(header_y),
                                        Inches(cw), Inches(0.5))
            rect.fill.solid()
            rect.fill.fore_color.rgb = NAVY
            rect.line.fill.background()
            add_text(s6, hdr, cx + 0.05, header_y + 0.07, cw - 0.1, 0.36,
                     size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
            cx += cw

        # Data rows
        for r_idx, (theme, row) in enumerate(pivot.iterrows()):
            ry = header_y + 0.5 + r_idx * 0.72
            bg = WHITE if r_idx % 2 == 0 else RGBColor(0xf8, 0xf9, 0xfa)
            cx = start_x
            for j, col in enumerate(col_headers):
                cw = col_widths[j]
                rect = s6.shapes.add_shape(1, Inches(cx), Inches(ry),
                                            Inches(cw), Inches(0.65))
                rect.fill.solid()
                rect.fill.fore_color.rgb = bg
                rect.line.color.rgb = RGBColor(0xe9, 0xec, 0xef)
                rect.line.width = Emu(6350)

                if col == "Theme":
                    val = str(theme)
                    tc_ = NAVY
                elif col == "Total":
                    val = str(int(row["Total"]))
                    tc_ = NAVY
                else:
                    sc = col.lower().replace("neutral/mixed", "neutral/mixed")
                    # match original column name
                    matched = next((c for c in scols if c.lower() == sc), None)
                    if matched and matched in row:
                        v   = int(row.get(matched, 0))
                        pct = v / row["Total"] * 100 if row["Total"] else 0
                        val = f"{v}  ({pct:.0f}%)"
                        if matched == "positive":   tc_ = GREEN
                        elif matched == "negative": tc_ = RED
                        else:                       tc_ = GRAY
                    else:
                        val = "—"
                        tc_ = GRAY

                add_text(s6, val, cx + 0.08, ry + 0.12, cw - 0.16, 0.45,
                         size=10, bold=(col == "Theme"), color=tc_, align=PP_ALIGN.LEFT)
                cx += cw

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Positive vs Negative Theme Comparison
    # ══════════════════════════════════════════════════════════════════════════
    if has_themes:
        s7 = prs.slides.add_slide(BLANK)
        slide_header(s7, "Positive vs Negative — Top Themes",
                     "What drives satisfaction vs complaints?")
        add_rect(s7, 0, 1.1, 13.33, 6.4, LIGHT)

        pos_exp = df[df["predicted_sentiment"] == "positive"]["themes"]\
                  .str.split(r",\s*").explode().str.strip()
        pos_exp = pos_exp[~pos_exp.isin(["FAILED","","NOT PROCESSED"])]
        pos_tc  = pos_exp.value_counts().head(6).reset_index()
        pos_tc.columns = ["Theme", "Count"]

        neg_exp = df[df["predicted_sentiment"] == "negative"]["themes"]\
                  .str.split(r",\s*").explode().str.strip()
        neg_exp = neg_exp[~neg_exp.isin(["FAILED","","NOT PROCESSED"])]
        neg_tc  = neg_exp.value_counts().head(6).reset_index()
        neg_tc.columns = ["Theme", "Count"]

        fig_pos = px.bar(pos_tc, x="Count", y="Theme", orientation="h",
                         color_discrete_sequence=["16a34a"],
                         labels={"Count": "Mentions"}, title="Top themes — Positive")
        fig_pos.update_traces(marker_color="#16a34a", marker_line_color="white",
                              marker_line_width=1)
        fig_pos.update_layout(bargap=0.2, paper_bgcolor="white",
                               yaxis=dict(categoryorder="total ascending"),
                               margin=dict(l=10, r=10, t=40, b=10))

        fig_neg = px.bar(neg_tc, x="Count", y="Theme", orientation="h",
                         labels={"Count": "Mentions"}, title="Top themes — Negative")
        fig_neg.update_traces(marker_color="#dc2626", marker_line_color="white",
                              marker_line_width=1)
        fig_neg.update_layout(bargap=0.2, paper_bgcolor="white",
                               yaxis=dict(categoryorder="total ascending"),
                               margin=dict(l=10, r=10, t=40, b=10))

        add_chart_image(s7, fig_pos, 0.3, 1.3, 6.3, 5.8, width_px=640, height_px=580)
        add_chart_image(s7, fig_neg, 6.8, 1.3, 6.3, 5.8, width_px=640, height_px=580)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Key Findings
    # ══════════════════════════════════════════════════════════════════════════
    s8 = prs.slides.add_slide(BLANK)
    slide_header(s8, "Key Findings & Recommendations", "")
    add_rect(s8, 0, 1.1, 13.33, 6.4, LIGHT)

    # Auto-generate findings from the data
    findings = []
    if pos / total >= 0.6:
        findings.append(f"Strong positive sentiment — {pos/total*100:.0f}% of reviews are positive, indicating overall customer satisfaction.")
    elif neg / total >= 0.4:
        findings.append(f"High negative sentiment — {neg/total*100:.0f}% of reviews are negative, requiring immediate attention.")
    else:
        findings.append(f"Mixed sentiment — {pos/total*100:.0f}% positive and {neg/total*100:.0f}% negative reviews detected.")

    if avg_conf and avg_conf < 70:
        findings.append(f"Model confidence is moderate ({avg_conf:.1f}%) — consider reviewing low-confidence predictions manually.")
    elif avg_conf:
        findings.append(f"Model confidence is high ({avg_conf:.1f}%) — predictions are reliable.")

    if has_themes:
        top_theme = exp_all.value_counts().idxmax() if not exp_all.empty else "N/A"
        findings.append(f"Most discussed topic is '{top_theme}' — ensure this area meets customer expectations.")

        if has_themes and not neg_tc.empty:
            top_neg_theme = neg_tc.iloc[0]["Theme"]
            findings.append(f"'{top_neg_theme}' drives the most negative feedback — prioritise improvements here.")

    if mixed / total >= 0.1:
        findings.append(f"{mixed/total*100:.0f}% of reviews are neutral or mixed — these customers are on the fence and could be converted with targeted improvements.")

    recs = [
        "Share this report with operations teams to address the top negative theme.",
        "Monitor monthly sentiment trends to detect early warning signs.",
        "Use the Theme Extraction tab to drill into specific customer pain points.",
        "Re-run analysis after implementing changes to measure improvement.",
    ]

    # Left column — findings
    add_rect(s8, 0.4, 1.3, 5.9, 5.8, WHITE)
    add_text(s8, "Findings", 0.6, 1.4, 5.5, 0.5, size=15, bold=True, color=NAVY)
    for i, f in enumerate(findings[:5]):
        fy = 1.95 + i * 1.0
        dot = s8.shapes.add_shape(1, Inches(0.6), Inches(fy + 0.07),
                                   Inches(0.12), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s8, f, 0.85, fy, 5.2, 0.85, size=11, color=GRAY)

    # Right column — recommendations
    add_rect(s8, 6.9, 1.3, 5.9, 5.8, WHITE)
    add_text(s8, "Recommendations", 7.1, 1.4, 5.5, 0.5, size=15, bold=True, color=NAVY)
    for i, r in enumerate(recs):
        ry = 1.95 + i * 1.15
        num_box = s8.shapes.add_shape(1, Inches(7.1), Inches(ry),
                                       Inches(0.35), Inches(0.35))
        num_box.fill.solid(); num_box.fill.fore_color.rgb = NAVY
        num_box.line.fill.background()
        add_text(s8, str(i+1), 7.1, ry, 0.35, 0.35,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s8, r, 7.58, ry, 4.9, 0.85, size=11, color=GRAY)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Thank You / End
    # ══════════════════════════════════════════════════════════════════════════
    s9 = prs.slides.add_slide(BLANK)
    add_rect(s9, 0, 0, 13.33, 7.5, NAVY)
    add_rect(s9, 0, 5.2, 13.33, 2.3, ACCENT)
    add_text(s9, "Thank You", 0.8, 1.8, 11, 1.4,
             size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s9, "Customer Feedback Intelligence Platform",
             0.8, 3.3, 11, 0.6, size=18, color=RGBColor(0xCA, 0xDC, 0xFC))
    add_text(s9, f"{total:,} reviews  |  {datetime.now().strftime('%B %Y')}",
             0.8, 3.95, 11, 0.45, size=13, color=RGBColor(0x9F, 0xB3, 0xD8))
    add_text(s9, "Generated automatically — results should be reviewed alongside raw data.",
             0.8, 5.5, 11, 0.5, size=12, color=RGBColor(0xCA, 0xDC, 0xFC))

    # ── Save to bytes ──────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Shared CSS injected once ───────────────────────────────────────────────────
METRIC_CSS = """
<style>
[data-testid="stMetric"] {
    background-color: #f8f9fa;
    border: 1px solid #e9ecef;
    border-radius: 8px;
    padding: 16px;
    min-height: 110px;
}
[data-testid="stMetricLabel"] {
    font-size: 13px;
    font-weight: 600;
    color: #6b7280;
}
[data-testid="stMetricValue"] {
    font-size: 28px;
    font-weight: 700;
}
[data-testid="stMetricDelta"] {
    font-size: 12px;
    color: #6b7280 !important;
}
[data-testid="stMetricDelta"] svg {
    display: none;
}
</style>
"""


# ── Shared axis/layout helpers — keep every chart consistent ─────────────────

def _vertical_bar_axis():
    """Vertical bar / histogram: x=0–100% every 10%, y=whole integers from 0.
    Bars align exactly under their percentage label by using fixed bin edges."""
    return dict(
        xaxis=dict(
            tickformat=".0%",
            range=[-0.05, 1.05],          # small padding so end bars aren't clipped
            tickvals=[i / 10 for i in range(0, 11)],
            ticktext=[f"{i*10}%" for i in range(0, 11)],
            title_standoff=10,
        ),
        yaxis=dict(rangemode="tozero", dtick=1, tickformat="d"),
    )

def _horizontal_bar_axis(max_val):
    """Horizontal bar: x=whole integers, bars aligned to their tick via tick0=0, dtick.
    No decimals, no 0.5 steps."""
    nice_max = max(1, int(max_val) + 1)
    step = max(1, int(nice_max / 8))
    return dict(
        xaxis=dict(
            range=[0, nice_max + step * 0.5],  # padding so rightmost bar isn't clipped
            tick0=0,
            dtick=step,
            tickformat="d",
            rangemode="tozero",
        ),
        yaxis=dict(categoryorder="total ascending"),
    )

def _trend_bar_axis(date_vals):
    """Monthly trend bar: every bar labelled 'Mon YYYY', angled -45°, y=whole integers."""
    return dict(
        xaxis=dict(
            tickformat="%b %Y",
            tickmode="array",
            tickvals=list(date_vals),
            tickangle=-45,
            title_standoff=10,
        ),
        yaxis=dict(rangemode="tozero", dtick=1, tickformat="d"),
    )

def _base_layout(**extra):
    """White background, standard margins, applied to every chart."""
    layout = dict(
        paper_bgcolor="white",
        plot_bgcolor="white",
        margin=dict(l=10, r=20, t=45, b=70),
        bargap=0.2,       # consistent gap between bars across all charts
        bargroupgap=0.05, # gap between groups (for grouped bar charts)
    )
    layout.update(extra)
    return layout


def _chart_download(fig, filename, label="Download chart as PNG"):
    """Render a download button that saves the given Plotly figure as a PNG."""
    try:
        png_bytes = fig.to_image(format="png", scale=2)
        st.download_button(
            label=label,
            data=png_bytes,
            file_name=filename,
            mime="image/png",
        )
    except Exception:
        st.caption("Install kaleido to enable chart downloads: pip install kaleido")


def _pdf_download_button(df, key_suffix=""):
    """Render PDF and PowerPoint report download buttons."""
    st.markdown("---")
    st.markdown("**Download full report**")
    custom_title = st.text_input(
        "Report title",
        value="Feedback Intelligence Report",
        key=f"report_title_{key_suffix}",
    )

    col_pdf, col_pptx = st.columns(2)

    with col_pdf:
        if st.button("Generate PDF Report", key=f"gen_pdf_{key_suffix}"):
            with st.spinner("Building PDF report..."):
                try:
                    pdf_bytes = generate_pdf_report(df, report_title=custom_title)
                    st.download_button(
                        label="Download PDF Report",
                        data=pdf_bytes,
                        file_name=f"feedback_report_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                        mime="application/pdf",
                        key=f"dl_pdf_{key_suffix}",
                    )
                except Exception as e:
                    st.error(f"PDF generation failed: {e}")
                    st.info("Run: pip install kaleido")

    with col_pptx:
        if st.button("Generate PowerPoint", key=f"gen_pptx_{key_suffix}"):
            with st.spinner("Building PowerPoint deck..."):
                try:
                    pptx_bytes = generate_pptx(df, report_title=custom_title)
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_bytes,
                        file_name=f"feedback_deck_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key=f"dl_pptx_{key_suffix}",
                    )
                except Exception as e:
                    st.error(f"PowerPoint generation failed: {e}")
                    st.info("Run: pip install python-pptx kaleido")


def page_overview(df):
    # Inject metric card CSS so all cards sit at equal height
    st.markdown(METRIC_CSS, unsafe_allow_html=True)

    total    = len(df)
    pos      = int((df["predicted_sentiment"] == "positive").sum())
    neg      = int((df["predicted_sentiment"] == "negative").sum())
    mixed    = int(df["predicted_sentiment"].isin(["neutral","neutral/mixed"]).sum())
    avg_conf = df["confidence"].mean() if "confidence" in df.columns else None

    # KPI cards — all deltas use the same label format so height is consistent
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Total Reviews",     f"{total:,}",                    "100% of reviews")
    c2.metric("Positive",          f"{pos:,}",                      f"{pos/total*100:.1f}% of reviews")
    c3.metric("Negative",          f"{neg:,}",                      f"{neg/total*100:.1f}% of reviews")
    c4.metric("Neutral / Mixed",   f"{mixed:,}",                    f"{mixed/total*100:.1f}% of reviews")
    c5.metric("Avg Confidence",    f"{avg_conf*100:.1f}%" if avg_conf else "N/A", "model certainty")

    st.markdown("---")
    col_a, col_b = st.columns(2)

    with col_a:
        counts = df["predicted_sentiment"].value_counts().reset_index()
        counts.columns = ["sentiment", "count"]
        fig_pie = px.pie(counts, names="sentiment", values="count",
                         title="Predicted Sentiment Distribution",
                         color="sentiment", color_discrete_map=COLOUR_MAP)
        st.plotly_chart(fig_pie, use_container_width=True)
        _chart_download(fig_pie, "sentiment_distribution.png", "Download sentiment pie chart")

    with col_b:
        fig_trend = None
        if "date" in df.columns:
            try:
                df["date"] = pd.to_datetime(df["date"])
                td = df.groupby([pd.Grouper(key="date", freq="ME"), "predicted_sentiment"])\
                       .size().reset_index(name="count")
                if len(td) > 1:
                    fig_trend = px.line(td, x="date", y="count", color="predicted_sentiment",
                                        title="Sentiment Over Time (Monthly)",
                                        color_discrete_map=COLOUR_MAP)
                    st.plotly_chart(fig_trend, use_container_width=True)
                    _chart_download(fig_trend, "sentiment_trend.png", "Download trend chart")
                else:
                    st.info("Not enough date range for a trend chart.")
            except Exception:
                st.info("Date column could not be parsed.")
        else:
            st.info("No 'date' column found — add one to your CSV to see trends.")

    # ── Star rating breakdown ──────────────────────────────────────────────────
    if "stars" in df.columns:
        st.markdown("---")
        st.subheader("Star Rating vs Model Sentiment")
        st.write(
            "This shows how many reviews have each star rating, coloured by what "
            "the model predicted. Ideally 4-5 stars should be mostly green (positive) "
            "and 1-2 stars mostly red (negative). Mismatches highlight where the model "
            "disagreed with the reviewer's own rating."
        )

        try:
            star_df = df.copy()
            star_df["stars"] = pd.to_numeric(star_df["stars"], errors="coerce")
            star_df = star_df.dropna(subset=["stars"])
            star_df["stars"] = star_df["stars"].astype(int)
            star_df = star_df[star_df["stars"].between(1, 5)]

            # Bar chart — star rating on x, coloured by predicted sentiment
            star_counts = (
                star_df.groupby(["stars", "predicted_sentiment"])
                .size()
                .reset_index(name="count")
            )
            fig_stars = px.bar(
                star_counts,
                x="stars",
                y="count",
                color="predicted_sentiment",
                title="Star Rating Distribution — coloured by Model Prediction",
                color_discrete_map=COLOUR_MAP,
                barmode="stack",
                labels={"stars": "Star Rating (1=worst, 5=best)", "count": "Number of Reviews"},
            )
            fig_stars.update_traces(marker_line_color="white", marker_line_width=1)
            fig_stars.update_layout(**_base_layout(
                bargap=0.2,
                xaxis=dict(
                    tickmode="array",
                    tickvals=[1, 2, 3, 4, 5],
                    ticktext=["1 star", "2 stars", "3 stars", "4 stars", "5 stars"],
                ),
                yaxis=dict(rangemode="tozero", dtick=1, tickformat="d"),
                legend_title_text="Model Prediction",
            ))
            st.plotly_chart(fig_stars, use_container_width=True)
            _chart_download(fig_stars, "star_rating_breakdown.png", "Download star rating chart")

            # Agreement table — shows % model agreed with star rating
            st.markdown("**Agreement between star rating and model prediction**")
            st.caption(
                "Positive = 4-5 stars, Negative = 1-2 stars, Neutral = 3 stars. "
                "Agreement means the model prediction matches what the star rating implies."
            )

            def expected_sentiment(stars):
                if stars >= 4:   return "positive"
                elif stars <= 2: return "negative"
                else:            return "neutral/mixed"

            star_df["expected"] = star_df["stars"].apply(expected_sentiment)
            star_df["model_agrees"] = (
                star_df["predicted_sentiment"] == star_df["expected"]
            ) | (
                (star_df["predicted_sentiment"] == "neutral") &
                (star_df["expected"] == "neutral/mixed")
            )

            agree_by_star = (
                star_df.groupby("stars")["model_agrees"]
                .agg(["sum", "count"])
                .reset_index()
            )
            agree_by_star.columns = ["Stars", "Agreed", "Total"]
            agree_by_star["Agreement %"] = (
                agree_by_star["Agreed"] / agree_by_star["Total"] * 100
            ).round(1).astype(str) + "%"
            agree_by_star["Stars"] = agree_by_star["Stars"].astype(str) + " star"

            overall_agree = star_df["model_agrees"].mean() * 100
            st.dataframe(agree_by_star[["Stars","Agreed","Total","Agreement %"]],
                         use_container_width=True)
            st.caption(f"Overall model agreement with star ratings: **{overall_agree:.1f}%**")

        except Exception as e:
            st.info(f"Could not render star rating breakdown: {e}")

    st.markdown("---")
    st.markdown("**Data preview — first 10 rows**")
    preview_cols = [c for c in ["predicted_sentiment", "confidence", "is_mixed", "themes"]
                    if c in df.columns]
    st.dataframe(df[preview_cols].head(10), use_container_width=True)

    _pdf_download_button(df, key_suffix="overview")


def page_positive(df):
    pos_df = df[df["predicted_sentiment"] == "positive"].copy()
    st.markdown(f"### Positive Reviews — {len(pos_df):,} total")
    if pos_df.empty:
        st.info("No positive reviews found.")
        return

    if "confidence" in pos_df.columns:
        fig_hist = px.histogram(pos_df, x="confidence",
                                title="Confidence Score Distribution — Positive Reviews",
                                color_discrete_sequence=["#16a34a"])
        fig_hist.update_traces(
            marker_line_color="white", marker_line_width=2,
            xbins=dict(start=0.0, end=1.0, size=0.10),
        )
        fig_hist.update_layout(**_base_layout(
            xaxis_title="Model Confidence Score (higher = more certain)",
            yaxis_title="Number of Reviews",
            **_vertical_bar_axis(),
        ))
        st.plotly_chart(fig_hist, use_container_width=True)
        _chart_download(fig_hist, "positive_confidence.png", "Download confidence chart")

    if "themes" in pos_df.columns:
        st.markdown("**Most mentioned themes in positive reviews**")
        exp = pos_df["themes"].str.split(r",\s*").explode().str.strip()
        exp = exp[~exp.isin(["FAILED", "", "NOT PROCESSED"])]
        tc  = exp.value_counts().head(8).reset_index()
        tc.columns = ["Theme", "Count"]
        fig_bar = px.bar(tc, x="Count", y="Theme", orientation="h",
                         color_discrete_sequence=["#16a34a"])
        fig_bar.update_traces(marker_line_color="white", marker_line_width=1)
        fig_bar.update_layout(**_base_layout(
            xaxis_title="Number of reviews mentioning this theme",
            **_horizontal_bar_axis(tc["Count"].max()),
        ))
        st.plotly_chart(fig_bar, use_container_width=True)
        _chart_download(fig_bar, "positive_themes.png", "Download themes chart")

    st.markdown("**Sample positive reviews**")
    tc_name = _text_col(pos_df)
    hc, bc  = st.columns([3, 1])
    with bc:
        st.button("Refresh", key="pos_refresh")
    for _, row in pos_df.sample(min(10, len(pos_df))).iterrows():
        conf = f" *(confidence: {row['confidence']*100:.0f}%)*" if "confidence" in row else ""
        st.success(f'"{row[tc_name]}"{conf}')

    _pdf_download_button(df, key_suffix="positive")


def page_negative(df):
    neg_df = df[df["predicted_sentiment"] == "negative"].copy()
    st.markdown(f"### Negative Reviews — {len(neg_df):,} total")
    if neg_df.empty:
        st.info("No negative reviews found.")
        return

    if "confidence" in neg_df.columns:
        fig_hist = px.histogram(neg_df, x="confidence",
                                title="Confidence Score Distribution — Negative Reviews",
                                color_discrete_sequence=["#dc2626"])
        fig_hist.update_traces(
            marker_line_color="white", marker_line_width=2,
            xbins=dict(start=0.0, end=1.0, size=0.10),
        )
        fig_hist.update_layout(**_base_layout(
            xaxis_title="Model Confidence Score (higher = more certain)",
            yaxis_title="Number of Reviews",
            **_vertical_bar_axis(),
        ))
        st.plotly_chart(fig_hist, use_container_width=True)
        _chart_download(fig_hist, "negative_confidence.png", "Download confidence chart")

    if "themes" in neg_df.columns:
        st.markdown("**Most mentioned themes in negative reviews**")
        exp = neg_df["themes"].str.split(r",\s*").explode().str.strip()
        exp = exp[~exp.isin(["FAILED", "", "NOT PROCESSED"])]
        tc  = exp.value_counts().head(8).reset_index()
        tc.columns = ["Theme", "Count"]
        fig_bar = px.bar(tc, x="Count", y="Theme", orientation="h",
                         color_discrete_sequence=["#dc2626"])
        fig_bar.update_traces(marker_line_color="white", marker_line_width=1)
        fig_bar.update_layout(**_base_layout(
            xaxis_title="Number of reviews mentioning this theme",
            **_horizontal_bar_axis(tc["Count"].max()),
        ))
        st.plotly_chart(fig_bar, use_container_width=True)
        _chart_download(fig_bar, "negative_themes.png", "Download themes chart")

    st.markdown("**Sample negative reviews**")
    tc_name = _text_col(neg_df)
    hc, bc  = st.columns([3, 1])
    with bc:
        st.button("Refresh", key="neg_refresh")
    for _, row in neg_df.sample(min(10, len(neg_df))).iterrows():
        conf = f" *(confidence: {row['confidence']*100:.0f}%)*" if "confidence" in row else ""
        st.error(f'"{row[tc_name]}"{conf}')

    _pdf_download_button(df, key_suffix="negative")


def page_neutral(df):
    neu_df = df[df["predicted_sentiment"].isin(["neutral", "neutral/mixed"])].copy()
    st.markdown(f"### Neutral / Mixed Reviews — {len(neu_df):,} total")
    st.write(
        "These reviews sit in the middle — the model did not detect a clearly positive "
        "or negative tone. Some may be genuinely balanced, others may be vague or ambiguous."
    )
    if neu_df.empty:
        st.info("No neutral or mixed reviews found in this dataset.")
        return

    if "confidence" in neu_df.columns:
        fig_hist = px.histogram(neu_df, x="confidence",
                                title="Confidence Score Distribution — Neutral / Mixed Reviews",
                                color_discrete_sequence=["#6b7280"])
        fig_hist.update_traces(
            marker_line_color="white", marker_line_width=2,
            xbins=dict(start=0.0, end=1.0, size=0.10),
        )
        fig_hist.update_layout(**_base_layout(
            xaxis_title="Model Confidence Score (higher = more certain it is neutral)",
            yaxis_title="Number of Reviews",
            **_vertical_bar_axis(),
        ))
        st.plotly_chart(fig_hist, use_container_width=True)
        _chart_download(fig_hist, "neutral_confidence.png", "Download confidence chart")

    # Split into neutral vs mixed — metrics shown below histogram
    if "is_mixed" in neu_df.columns:
        truly_neutral = neu_df[neu_df["is_mixed"] == False]
        mixed_signal  = neu_df[neu_df["is_mixed"] == True]
        col_a, col_b  = st.columns(2)
        with col_a:
            st.metric("Purely Neutral", f"{len(truly_neutral):,}",
                      "No strong positive or negative language")
        with col_b:
            st.metric("Mixed Signal", f"{len(mixed_signal):,}",
                      "Contains both positive and negative language")

    if "themes" in neu_df.columns:
        st.markdown("**Most mentioned themes in neutral / mixed reviews**")
        exp = neu_df["themes"].str.split(r",\s*").explode().str.strip()
        exp = exp[~exp.isin(["FAILED", "", "NOT PROCESSED"])]
        tc  = exp.value_counts().head(8).reset_index()
        tc.columns = ["Theme", "Count"]
        fig_bar = px.bar(tc, x="Count", y="Theme", orientation="h",
                         color_discrete_sequence=["#6b7280"])
        fig_bar.update_traces(marker_line_color="white", marker_line_width=1)
        fig_bar.update_layout(**_base_layout(
            xaxis_title="Number of reviews mentioning this theme",
            **_horizontal_bar_axis(tc["Count"].max()),
        ))
        st.plotly_chart(fig_bar, use_container_width=True)
        _chart_download(fig_bar, "neutral_themes.png", "Download themes chart")

    st.markdown("**Sample neutral / mixed reviews**")
    tc_name = _text_col(neu_df)
    hc, bc  = st.columns([3, 1])
    with bc:
        st.button("Refresh", key="neu_refresh")
    for _, row in neu_df.sample(min(10, len(neu_df))).iterrows():
        conf  = f" *(confidence: {row['confidence']*100:.0f}%)*" if "confidence" in row else ""
        label = "Mixed" if ("is_mixed" in row and row["is_mixed"]) else "Neutral"
        st.warning(f'**[{label}]** "{row[tc_name]}"{conf}')

    _pdf_download_button(df, key_suffix="neutral")


def page_themes(df):
    if "themes" not in df.columns:
        st.warning("No themes column found. Run the analysis pipeline first.")
        return

    exp_all = df["themes"].str.split(r",\s*").explode().str.strip()
    exp_all = exp_all[~exp_all.isin(["FAILED", "", "NOT PROCESSED"])]
    theme_counts = exp_all.value_counts().head(20).reset_index()
    theme_counts.columns = ["Theme", "Count"]

    # Bar chart first, full width
    fig_theme_bar = px.bar(theme_counts, x="Count", y="Theme", orientation="h",
                           title="Most Common Themes", color="Count",
                           color_continuous_scale="Viridis")
    fig_theme_bar.update_traces(marker_line_color="white", marker_line_width=1)
    fig_theme_bar.update_layout(**_base_layout(
        xaxis_title="Number of mentions",
        coloraxis_showscale=False,
        **_horizontal_bar_axis(theme_counts["Count"].max()),
    ))
    st.plotly_chart(fig_theme_bar, use_container_width=True)
    _chart_download(fig_theme_bar, "top_themes.png", "Download themes chart")

    # Table below chart
    st.markdown("**Theme mention counts**")
    st.dataframe(theme_counts, use_container_width=True)

    st.markdown("---")
    st.markdown("**Sentiment Breakdown by Theme**")
    df_exp = df.assign(Theme=df["themes"].str.split(r",\s*")).explode("Theme")
    df_exp["Theme"] = df_exp["Theme"].str.strip()
    df_exp = df_exp[~df_exp["Theme"].isin(["FAILED", "", "NOT PROCESSED"])].reset_index(drop=True)

    if not df_exp.empty:
        pivot = pd.crosstab(df_exp["Theme"], df_exp["predicted_sentiment"],
                            margins=True, margins_name="Total")
        scols = [c for c in ["positive", "negative", "neutral", "neutral/mixed"] if c in pivot.columns]
        for c in scols:
            pivot[c+" (%)"] = (pivot[c] / pivot["Total"] * 100).round(1)
        ordered = []
        for c in scols:
            ordered.extend([c, c+" (%)"])
        ordered.append("Total")
        st.dataframe(pivot[ordered], use_container_width=True)

        st.markdown("---")
        st.markdown("**Interactive Theme Sentiment Distribution**")
        unique_themes = sorted(df_exp["Theme"].unique().tolist())
        selected      = st.selectbox("Select a theme:", unique_themes, key="theme_select")
        t_data        = df_exp[df_exp["Theme"] == selected]
        t_counts      = t_data["predicted_sentiment"].value_counts().reset_index()
        t_counts.columns = ["sentiment", "count"]

        # Pie chart full width
        fig_pie = px.pie(t_counts, names="sentiment", values="count",
                         title=f"Sentiment split for '{selected}'",
                         color="sentiment", color_discrete_map=COLOUR_MAP)
        st.plotly_chart(fig_pie, use_container_width=True)
        _chart_download(fig_pie, f"theme_{selected.replace(' ','_')}_sentiment.png",
                        "Download theme sentiment chart")

        # Sentiment counts table below pie
        st.markdown(f"**{selected} — review counts by sentiment**")
        st.dataframe(t_counts.rename(columns={"sentiment":"Sentiment","count":"Reviews"}),
                     use_container_width=True)

        # Top phrases
        st.markdown("---")
        st.markdown(f"**Top phrases in '{selected}'**")
        try:
            sent_docs = df_exp.groupby("Theme")["clean_text"].apply(
                lambda x: " ".join(x.dropna())).to_dict()
            if selected in sent_docs and sent_docs[selected].strip():
                extra_stop = {"review","user","star","stars","https","http","amp","just","like","im"}
                sw  = list(set(sk_text.ENGLISH_STOP_WORDS) | extra_stop)
                tv  = TfidfVectorizer(ngram_range=(2, 3), stop_words=sw)
                mat = tv.fit_transform(list(sent_docs.values()))
                idx = list(sent_docs.keys()).index(selected)
                dense  = mat.toarray()
                scores = dense[idx] * (dense.argmax(axis=0) == idx)
                top_i  = scores.argsort()[-10:][::-1]
                phrases = [tv.get_feature_names_out()[i] for i in top_i if scores[i] > 0]
                pcounts = [sent_docs[selected].count(p) for p in phrases]
                pf = pd.DataFrame({"Phrase": phrases, "Count": pcounts})\
                       .sort_values("Count", ascending=True)
                fig3 = px.bar(pf, x="Count", y="Phrase", orientation="h")
                fig3.update_traces(marker_line_color="white", marker_line_width=1)
                fig3.update_layout(**_base_layout(
                    xaxis_title="Mentions",
                    yaxis_title="",
                    **_horizontal_bar_axis(pf["Count"].max()),
                ))
                st.plotly_chart(fig3, use_container_width=True)
        except Exception as e:
            st.info(f"Could not extract phrases: {e}")

        # Time trends — scrollable chart + table below
        if "date" in df.columns:
            st.markdown("---")
            st.subheader("Time-Based & Emergent Trends")
            try:
                theme_time = (
                    df_exp.groupby([pd.Grouper(key="date", freq="ME"), "Theme"])
                    .size().reset_index(name="count")
                )
                if len(theme_time["date"].unique()) > 1:
                    st.markdown(f"**Monthly volume for '{selected}'**")
                    sel_tt = theme_time[theme_time["Theme"] == selected]

                    # Dynamically size width so bars have room — minimum 800px, 20px per bar
                    n_bars    = len(sel_tt)
                    chart_w   = max(800, n_bars * 28)
                    bar_h     = 420

                    fig_t = px.bar(sel_tt, x="date", y="count",
                                   labels={"date": "Month", "count": "Mentions"})
                    fig_t.update_traces(marker_line_color="white", marker_line_width=1,
                                        marker_color="#065a82")
                    fig_t.update_layout(
                        width=chart_w,
                        height=bar_h,
                        **_base_layout(
                            margin=dict(l=10, r=20, t=45, b=120),
                            **_trend_bar_axis(sel_tt["date"]),
                        ),
                    )
                    # Render inside a scrollable container
                    st.markdown(
                        f'<div style="overflow-x:auto; width:100%;">'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                    st.plotly_chart(fig_t, use_container_width=False)

                    # Emergent themes table — below the chart
                    st.markdown("**Emergent Themes (month-over-month)**")
                    st.write("Change in mention volume between the two most recent months.")
                    months     = sorted(theme_time["date"].unique())
                    curr_month = months[-1]
                    prev_month = months[-2]
                    curr_df = theme_time[theme_time["date"]==curr_month].set_index("Theme")
                    prev_df = theme_time[theme_time["date"]==prev_month].set_index("Theme")
                    emer = curr_df[["count"]].join(prev_df[["count"]],
                                                    lsuffix="_curr", rsuffix="_prev",
                                                    how="outer").fillna(0)
                    emer["Change"] = (emer["count_curr"] - emer["count_prev"]).astype(int)
                    emer["count_curr"] = emer["count_curr"].astype(int)
                    emer["count_prev"] = emer["count_prev"].astype(int)
                    st.dataframe(
                        emer.sort_values("Change", ascending=False)
                        [["count_prev","count_curr","Change"]]
                        .rename(columns={"count_prev":"Prev Month",
                                         "count_curr":"Current Month",
                                         "Change":"Momentum"}),
                        use_container_width=True)
                else:
                    st.info("Need at least 2 months of data for trend momentum.")
            except Exception as e:
                st.warning(f"Could not calculate trends: {e}")

        # Deep dive verbatims
        st.markdown("---")
        st.markdown("**Read actual reviews by theme + sentiment**")
        dd_col1, dd_col2 = st.columns(2)
        with dd_col1:
            dd_theme = st.selectbox("Theme:", unique_themes, key="dd_theme")
        with dd_col2:
            dd_sent  = st.selectbox("Sentiment:", ["negative","positive","neutral","neutral/mixed"], key="dd_sent")
        dd_data = df_exp[(df_exp["Theme"]==dd_theme) & (df_exp["predicted_sentiment"]==dd_sent)]
        if dd_data.empty:
            st.info(f"No {dd_sent} reviews found for '{dd_theme}'.")
        else:
            hc, bc = st.columns([3, 1])
            with bc:
                st.button("Refresh", key="dd_refresh")
            tc_name = _text_col(dd_data)
            for rev in dd_data.sample(min(5, len(dd_data)))[tc_name].tolist():
                st.info(f'"{rev}"')

    _pdf_download_button(df, key_suffix="themes")


def page_outliers(df):
    st.markdown("### Outlier & Edge-Case Reviews")
    st.write(
        "**Low-confidence predictions** are reviews the model was unsure about — "
        "they scored close to 50/50 between two sentiments. "
        "**Mixed-signal reviews** contain both positive and negative language in the same review."
    )

    if "confidence" in df.columns:
        st.markdown("---")
        st.markdown("#### Low-Confidence Predictions")
        threshold = st.slider(
            "Show reviews below this confidence level:",
            0.40, 0.90, 0.60, 0.05,
            help="Lower = model was very unsure. 60% means the model was only 60% sure of its prediction."
        )
        low_conf = df[df["confidence"] < threshold].copy()
        st.markdown(
            f"**{len(low_conf):,} reviews** were predicted with less than "
            f"**{threshold*100:.0f}% confidence** — these may be worth checking manually."
        )

        if not low_conf.empty:
            fig_out = px.histogram(
                low_conf, x="confidence", color="predicted_sentiment",
                title="How uncertain was the model? (grouped by predicted sentiment)",
                color_discrete_map=COLOUR_MAP,
                barmode="group",
            )
            fig_out.update_traces(
                marker_line_color="white", marker_line_width=2,
                xbins=dict(start=0.0, end=1.0, size=0.05),
            )
            fig_out.update_layout(**_base_layout(
                xaxis_title="Model Confidence Score (0.50 = 50% sure, almost a coin flip)",
                yaxis_title="Number of Reviews",
                legend_title_text="Predicted Sentiment",
                **_vertical_bar_axis(),
            ))
            st.plotly_chart(fig_out, use_container_width=True)
            _chart_download(fig_out, "low_confidence_distribution.png",
                            "Download low-confidence chart")

            tc_name   = _text_col(low_conf)
            show_cols = [tc_name, "predicted_sentiment", "confidence"] + \
                        [c for c in ["themes", "is_mixed"] if c in low_conf.columns]
            st.markdown("**Reviews sorted by lowest confidence first:**")
            st.dataframe(
                low_conf[show_cols].sort_values("confidence").head(30),
                use_container_width=True
            )
            st.download_button(
                "Download low-confidence reviews",
                data=low_conf.to_csv(index=False),
                file_name="low_confidence_reviews.csv",
                mime="text/csv",
            )

    st.markdown("---")

    if "is_mixed" in df.columns:
        st.markdown("#### Mixed-Signal Reviews")
        st.write(
            "These reviews contain contrast words like *'but', 'however', 'although'* "
            "alongside both positive and negative vocabulary — e.g. *'Great coffee but terrible service'*."
        )
        mixed_df = df[df["is_mixed"] == True].copy()
        st.markdown(f"**{len(mixed_df):,} mixed-signal reviews** detected")

        if not mixed_df.empty:
            tc_name   = _text_col(mixed_df)
            show_cols = [tc_name, "predicted_sentiment", "confidence"] + \
                        [c for c in ["themes"] if c in mixed_df.columns]
            st.dataframe(mixed_df[show_cols].head(30), use_container_width=True)
            st.download_button(
                "Download mixed-signal reviews",
                data=mixed_df.to_csv(index=False),
                file_name="mixed_signal_reviews.csv",
                mime="text/csv",
            )
        else:
            st.info("No mixed-signal reviews detected in this dataset.")

    _pdf_download_button(df, key_suffix="outliers")


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    st.title("Feedback Intelligence Platform")
    st.markdown(
        "Upload a review CSV in the sidebar, click **Run Analysis**, "
        "then navigate results using the tabs below."
    )

    model, vectorizer, train_accuracy = load_or_train_model()
    if model is None or vectorizer is None:
        st.error("Model could not be loaded or trained.")
        return
    if train_accuracy is not None:
        st.success(f"Model trained — accuracy: {train_accuracy:.4f}")

    # ── Sidebar ────────────────────────────────────────────────────────────────
    st.sidebar.header("Upload & Run")
    uploaded_file = st.sidebar.file_uploader(
        "Upload CSV for analysis", type="csv",
        help="CSV must have a 'text', 'clean_text', or 'raw_text' column.")

    st.sidebar.markdown("---")
    run_button = st.sidebar.button("Run Analysis", use_container_width=True)

    if st.sidebar.button("Reset / Clear Data", use_container_width=True):
        st.session_state.clear()
        st.rerun()

    # ── Run pipeline ───────────────────────────────────────────────────────────
    if uploaded_file and run_button:
        try:
            df = pd.read_csv(uploaded_file)
            df = preprocess_reviews(df)
            df = predict_reviews(df, model, vectorizer)
            df = extract_themes(df, THEMES)
            st.session_state.analyzed_df = df
        except Exception as exc:
            st.error(f"Error during analysis: {exc}")

    # ── Tabbed results ─────────────────────────────────────────────────────────
    if "analyzed_df" in st.session_state:
        df = st.session_state.analyzed_df
        st.markdown("---")

        tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
            "Overview",
            "Positive Reviews",
            "Neutral Reviews",
            "Negative Reviews",
            "Theme Extraction",
            "Outliers",
        ])

        with tab1:
            page_overview(df)
        with tab2:
            page_positive(df)
        with tab3:
            page_neutral(df)
        with tab4:
            page_negative(df)
        with tab5:
            page_themes(df)
        with tab6:
            page_outliers(df)

    elif not uploaded_file:
        st.info("Upload a CSV file in the sidebar and click Run Analysis to get started.")


if __name__ == "__main__":
    main()
